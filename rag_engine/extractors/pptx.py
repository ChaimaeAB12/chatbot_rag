from pptx import Presentation
from io import BytesIO

# Import de la fonction d'analyse d'image depuis image.py
from rag_engine.extractors.image import extract_text_from_image

def extract_text_from_pptx(file_bytes: bytes) -> str:
    prs = Presentation(BytesIO(file_bytes))
    full_content = []
    
    for idx, slide in enumerate(prs.slides, start=1):
        slide_text = [f"<<SLIDE {idx} START>>", f"🖼️ Slide {idx}"]

        # Titre
        if slide.shapes.title and slide.shapes.title.text:
            slide_text.append(f"📌 Titre : {slide.shapes.title.text.strip()}")

        # Texte et puces
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip():
                for para in shape.text_frame.paragraphs:
                    line = ""
                    if para.level == 0:
                        line += "• "
                    elif para.level == 1:
                        line += "    ◦ "
                    elif para.level == 2:
                        line += "        ▪ "
                    line += para.text.strip()
                    if line:
                        slide_text.append(line)

        # Tableaux
        for shape in slide.shapes:
            if shape.has_table:
                slide_text.append("📊 Tableau :")
                table = shape.table
                for row in table.rows:
                    row_text = " | ".join(cell.text.strip() for cell in row.cells)
                    slide_text.append(row_text)

        # Images
        for shape in slide.shapes:
            if shape.shape_type == 13:  # PICTURE
                image_bytes = shape.image.blob
                try:
                    image_analysis = extract_text_from_image(image_bytes)
                    slide_text.append(image_analysis)
                except Exception as e:
                    slide_text.append(f"⚠️ Erreur extraction image : {e}")

        # Notes du présentateur
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            note = slide.notes_slide.notes_text_frame.text.strip()
            if note:
                slide_text.append("📝 Notes du présentateur :")
                slide_text.append(note)

        slide_text.append(f"<<SLIDE {idx} END>>")
        full_content.append("\n".join(slide_text))

    return "\n\n===\n\n".join(full_content) if full_content else "Aucun contenu trouvé dans la présentation."
